from pinecone import Pinecone, ServerlessSpec
import pinecone
import hashlib

import fitz
import unstructured

from langchain.document_loaders import DirectoryLoader, PyPDFLoader, PyMuPDFLoader, UnstructuredWordDocumentLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings
from langchain.schema import Document

from openai import OpenAI

import time
import datetime
from datetime import datetime as dt
import os

import nltk
nltk.download('punkt')  # Pour la tokenisation
nltk.download('averaged_perceptron_tagger')  # Pour le POS tagging
nltk.download('averaged_perceptron_tagger_eng')
nltk.download('wordnet')  # Pour la lemmatisation

from pptx import Presentation


from transformers import AutoTokenizer
from sentence_transformers import SentenceTransformer

import requests
import json
from bs4 import BeautifulSoup
from newspaper import Article

from dateutil.relativedelta import relativedelta
from dateutil import parser
import re

from pathlib import Path




OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
PINECONE_INDEX_NAME = os.getenv("PINECONE_INDEX_NAME")

# Modèle
model_name = "multilingual-e5-large"
dimension_model = 1024
chunk_size = 6000
chunk_overlap = 0
max_tokens = 96
tokenizer = AutoTokenizer.from_pretrained("intfloat/multilingual-e5-large")

documents_path = "/content/drive/MyDrive/RAG_TEST_2"


# Initialisation de Pinecone (création si la bdd spécifiée n'existe pas)

def init_pinecone_db(PINECONE_INDEX_NAME,PINECONE_API_KEY, dimension_model):

  pc = Pinecone(
      api_key=PINECONE_API_KEY
  )
  if PINECONE_INDEX_NAME not in pc.list_indexes().names():

    index_name = PINECONE_INDEX_NAME

    pc.create_index(
        name=index_name,
        dimension=dimension_model,
        metric="cosine",
        spec=ServerlessSpec(
            cloud="aws",
            region="us-east-1"
            #cloud="gcp", --> Passer au plan Standard (https://docs.pinecone.io/troubleshooting/available-cloud-regions)
            #region="europe-west4"
        )
    )# ATTENTION --> BIEN CHOISIR LA REGION (SECURITE)
    time.sleep(20)  # Attente pour la création de l'index

  index = pc.describe_index(name=PINECONE_INDEX_NAME)
  print(index.name)
  return index


# === 🔹 Fonction pour générer des embeddings ===
def get_embedding(texts):
    """
    Génère des embeddings avec un modèle externe (ex : OpenAI, Ollama, Cohere)
    """

    pc = Pinecone(
      api_key=PINECONE_API_KEY
   )
    embeddings = pc.inference.embed(
        model=model_name,
        inputs=texts,
        parameters={"input_type": "passage", "truncate": "END"}
    )
    return [e["values"] for e in embeddings]  # Extraction des valeurs des embeddings

# === 🔹 Fonction pour nettoyer les pdf ===

def extract_clean_text_from_pdf(pdf_path):
    """Extrait et nettoie le texte d'un fichier PDF."""
    doc = fitz.open(pdf_path)
    text = ""

    for page in doc:
        text += page.get_text("text") + "\n"

    # Nettoyage du texte
    text = text.replace("\n", " ").replace("\n\t"," ").strip()  # Supprime les sauts de ligne inutiles
    text = " ".join(text.split())  # Supprime les espaces multiples
    return text

def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)


    text = extract_text_from_pptx("mon_presentation.pptx")
    return text



def internal_rag(documents_path, PINECONE_INDEX_NAME, chunk_size, chunk_overlap):
    """
    Charge et embedd les documents internes (.docx et .pdf) stockés dans un dossier.
    """
    pc = Pinecone(
      api_key=PINECONE_API_KEY
    )

    # Charger les documents .docx
    loader_docx = DirectoryLoader(documents_path, glob='**/*.txt')
    documents_docx = loader_docx.load()


    # Charger et nettoyer les documents .pdf
    pdf_files = [os.path.join(root, file)
                 for root, _, files in os.walk(documents_path)
                 for file in files if file.endswith('.pdf')]

    documents_pdf = []
    for pdf_file in pdf_files:
        cleaned_text = extract_clean_text_from_pdf(pdf_file)
        if cleaned_text:
            documents_pdf.append(Document(page_content=cleaned_text, metadata={"source": pdf_file}))

    # Charger les documents .pptx
    pptx_files = [os.path.join(root, file)
                 for root, _, files in os.walk(documents_path)
                 for file in files if file.endswith('.pptx')]
    documents_pptx = []
    for pptx_file in pptx_files:
        text = extract_text_from_pptx(pptx_file)
        if cleaned_text:
            documents_pptx.append(Document(page_content=cleaned_text, metadata={"source": pptx_file}))


    # Fusionner les documents
    documents = documents_docx + documents_pdf + documents_pptx

    print(f"✅ {len(documents)} documents chargés !")
    print(documents)
    if not documents:
        print("⚠️ Aucun document trouvé.")
        return

    # Fractionner les documents
    vectors = []  # Liste pour stocker tous les vecteurs

    for document in documents:
        document.page_content = document.page_content.strip()
        titre_docu = document.metadata
        titre_doc = Path(titre_docu['source']).stem
        print(titre_doc)

        text_splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)

        # On divise seulement le document actuel, pas la liste complète
        chunks = text_splitter.split_documents([document])

        # Séparer la chaîne en utilisant les chiffres comme séparateurs
        parties = re.split(r'\d+', titre_doc)
        # Supprimer les éventuelles chaînes vides
        parties = [p for p in parties if p]
        offre_doc = parties[0]  # 'un'
        service_doc = parties[1]  # 'deux'

        # Préparer les textes à embedder
        data = [{"id": f"DOC-TIBCO-{hashlib.md5(chunk.page_content.encode()).hexdigest()}",
                "text": chunk.page_content,
                "offre": offre_doc,
                "service": service_doc}
                for chunk in chunks]

        texts = [d['text'] for d in data]

        # Générer les embeddings
        embeddings = get_embedding(texts)

        # Ajouter les vecteurs créés à la liste
        vectors.extend([
            {"id": d["id"], "values": e, "metadata": {"type": "DOC-TIBCO", "text": d["text"], "offre": d['offre'], "service":d['service']}}
            for d, e in zip(data, embeddings)
        ])

    # Envoi des données accumulées vers Pinecone après la boucle
    if vectors:
        index_pc = pc.Index(PINECONE_INDEX_NAME)
        index_pc.upsert(vectors=vectors, namespace="DOC-TIBCO")
        print(f"✅ {len(vectors)} morceaux de documents indexés dans Pinecone !")

if __name__ == '__main__':
  init_pinecone_db(PINECONE_INDEX_NAME,PINECONE_API_KEY, dimension_model)
  internal_rag(documents_path, PINECONE_INDEX_NAME, chunk_size, chunk_overlap)